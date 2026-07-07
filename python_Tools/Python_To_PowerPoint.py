"""
================================================================================
脚本名称: 通用 Markdown 到 PPT 自动生成工具
功能描述:
    - 从 .md 文件读取内容，自动构建 PPT 幻灯片
    - 支持动态内容提取和智能排版
    - 一级标题(#) 作为幻灯片标题；二级标题(##) 作为小节标题
    - 项目符号(- * +) 和文本作为内容
    
使用示例:
    python Python_To_PowerPoint.py input.md output.pptx
    python Python_To_PowerPoint.py input.md
    python Python_To_PowerPoint.py  # 交互式输入

依赖库: pip install python-pptx

作者: Fan Zhen
最后修改日期: 2026-04-14
================================================================================
"""
import sys
import os
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- 定义专业配色方案 ---
COLOR_TEAL = RGBColor(15, 118, 110)      # 顶部栏主色
COLOR_WHITE = RGBColor(255, 255, 255)    # 顶部栏文字色
COLOR_TEXT = RGBColor(30, 41, 59)        # 正文文字色
COLOR_ACCENT = RGBColor(59, 130, 246)    # 强调色
COLOR_LIGHT_BG = RGBColor(248, 250, 252) # 浅色背景


class MarkdownParser:
    """Markdown 文档解析器：将 MD 内容提取为结构化的幻灯片数据"""
    
    def __init__(self, md_file_path):
        self.md_file_path = md_file_path
        self.slides_data = []
        self.parse()
    
    def parse(self):
        """从 Markdown 文件解析幻灯片数据"""
        with open(self.md_file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # 按一级标题分割，作为幻灯片
        slides_text = re.split(r'^# ', content, flags=re.MULTILINE)
        
        for slide_text in slides_text:
            if not slide_text.strip():
                continue
            
            lines = slide_text.strip().split('\n')
            slide_title = lines[0].strip()
            
            slide_data = {
                'title': slide_title,
                'sections': [],
                'bullet_points': []
            }
            
            current_section = None
            
            for line in lines[1:]:
                line = line.strip()
                
                if not line:
                    continue
                
                # 二级标题作为小节标题
                if line.startswith('## '):
                    if current_section is not None:
                        slide_data['sections'].append(current_section)
                    current_section = {
                        'title': line.replace('## ', '').strip(),
                        'points': []
                    }
                # 项目符号 (- 或 * 或 +)
                elif line.startswith(('- ', '* ', '+ ')):
                    point = re.sub(r'^[*\-+]\s*', '', line).strip()
                    if current_section is not None:
                        current_section['points'].append(point)
                    else:
                        slide_data['bullet_points'].append(point)
                # 普通文本作为项目符号
                elif line and not line.startswith('#'):
                    if current_section is not None:
                        current_section['points'].append(line)
                    else:
                        slide_data['bullet_points'].append(line)
            
            # 添加最后一个 section
            if current_section is not None:
                slide_data['sections'].append(current_section)
            
            self.slides_data.append(slide_data)


def create_styled_slide(prs, layout_idx, header_text):
    """创建带有专业深色 Header Bar 的幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # 柔和底色
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_LIGHT_BG
    bg.line.visible = False
    
    # 绘制顶部深色导航栏
    shapes = slide.shapes
    header_bar = shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLOR_TEAL
    header_bar.line.visible = False

    # 在 Header Bar 中添加标题
    title_box = shapes.add_textbox(
        Inches(0.5), Inches(0.15), prs.slide_width - Inches(1), Inches(0.5)
    )
    tf = title_box.text_frame
    tf.text = header_text
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.LEFT

    return slide


def add_content_to_slide(slide, bullet_points, slide_width, slide_height):
    """为幻灯片添加项目符号列表内容"""
    content_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.45),
        Inches(1.0),
        slide_width - Inches(0.9),
        slide_height - Inches(1.35),
    )
    content_card.fill.solid()
    content_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_card.line.color.rgb = RGBColor(226, 232, 240)

    body_shape = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.25), slide_width - Inches(1.6), slide_height - Inches(1.95)
    )
    tf = body_shape.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        p.text = point
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(8)
        p.space_after = Pt(6)
        p.line_spacing = 1.3
        p.level = 0


def add_content_with_sections(slide, sections, slide_width, slide_height):
    """为幻灯片添加多个分节的内容"""
    content_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.45),
        Inches(1.0),
        slide_width - Inches(0.9),
        slide_height - Inches(1.35),
    )
    content_card.fill.solid()
    content_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_card.line.color.rgb = RGBColor(226, 232, 240)

    body_shape = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.25), slide_width - Inches(1.6), slide_height - Inches(1.95)
    )
    tf = body_shape.text_frame
    tf.word_wrap = True
    
    first = True
    for section in sections:
        if not first:
            tf.add_paragraph()
        first = False
        
        # 小标题
        p_title = tf.paragraphs[0] if not first else tf.add_paragraph()
        p_title.text = section.get('title', '')
        p_title.font.name = "Microsoft YaHei"
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_ACCENT
        p_title.space_before = Pt(0)
        p_title.space_after = Pt(4)
        
        # 项目符号
        for point in section.get('points', []):
            p = tf.add_paragraph()
            p.text = point
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(15)
            p.font.color.rgb = COLOR_TEXT
            p.space_before = Pt(2)
            p.space_after = Pt(2)
            p.line_spacing = 1.2
            p.level = 1


def create_title_slide(prs, title_text, subtitle_text='', info_lines=None):
    """创建标题页幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 背景
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(241, 245, 249)
    title_bg.line.fill.background()

    # 顶部装饰条
    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.25)
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_TEAL
    accent_bar.line.fill.background()

    # 主标题
    title_box = slide.shapes.add_textbox(
        Inches(0.9), Inches(2.0), Inches(11.8), Inches(1.5)
    )
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_p.text = title_text
    title_p.font.name = "Microsoft YaHei"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_TEAL

    # 副标题
    if subtitle_text:
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.9), Inches(3.7), Inches(11.8), Inches(1.2)
        )
        subtitle_tf = subtitle_box.text_frame
        subtitle_p = subtitle_tf.paragraphs[0]
        subtitle_p.text = subtitle_text
        subtitle_p.font.name = "Microsoft YaHei"
        subtitle_p.font.size = Pt(24)
        subtitle_p.font.color.rgb = COLOR_ACCENT

    # 底部信息
    if info_lines:
        info_box = slide.shapes.add_textbox(
            Inches(0.9), Inches(5.2), Inches(11.8), Inches(1.8)
        )
        info_tf = info_box.text_frame
        info_tf.word_wrap = True
        
        for idx, line in enumerate(info_lines):
            p = info_tf.paragraphs[0] if idx == 0 else info_tf.add_paragraph()
            p.text = line
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(51, 65, 85)
            p.space_before = Pt(4)


def create_qa_slide(prs):
    """创建 Q&A 幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    qa_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    qa_bg.fill.solid()
    qa_bg.fill.fore_color.rgb = COLOR_LIGHT_BG
    qa_bg.line.fill.background()

    qa_accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.25)
    )
    qa_accent.fill.solid()
    qa_accent.fill.fore_color.rgb = COLOR_TEAL
    qa_accent.line.fill.background()

    # Q&A 内容
    qa_box = slide.shapes.add_textbox(
        Inches(1), Inches(2.5), Inches(11.333), Inches(3)
    )
    qa_tf = qa_box.text_frame
    qa_tf.word_wrap = True
    
    qa_lines = ['Q & A', '感谢聆听，欢迎提问！']
    for idx, line in enumerate(qa_lines):
        p = qa_tf.paragraphs[0] if idx == 0 else qa_tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(54 if idx == 0 else 32)
        p.font.bold = True if idx == 0 else False
        p.font.color.rgb = COLOR_TEAL if idx == 0 else RGBColor(71, 85, 105)
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0)
        p.space_after = Pt(8)


def create_presentation_from_markdown(md_file, output_file=None):
    """从 Markdown 文件创建 PPT"""
    
    # 验证输入文件
    if not os.path.exists(md_file):
        raise FileNotFoundError(f'Markdown 文件不存在: {md_file}')
    
    # 生成输出文件名
    if output_file is None:
        output_file = Path(md_file).stem + '.pptx'
    
    print(f'正在处理 Markdown 文件: {md_file}')
    
    # 解析 Markdown
    parser = MarkdownParser(md_file)
    slides_data = parser.slides_data
    
    if not slides_data:
        print('警告: 没有从 Markdown 文件中解析到幻灯片数据')
        return
    
    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    print(f'已解析 {len(slides_data)} 个幻灯片')
    
    # 第一个幻灯片作为标题页
    if slides_data:
        first_slide = slides_data[0]
        create_title_slide(
            prs, 
            first_slide['title'],
            subtitle_text='',
            info_lines=['日期: 2026年4月']
        )
        slides_data = slides_data[1:]  # 移除第一个，避免重复
    
    # 创建内容幻灯片
    for idx, slide_data in enumerate(slides_data):
        slide_title = slide_data['title']
        sections = slide_data['sections']
        bullet_points = slide_data['bullet_points']
        
        slide = create_styled_slide(prs, 1, f'Slide {idx + 2}: {slide_title}')
        
        # 优先使用 sections，次之使用 bullet_points
        if sections:
            add_content_with_sections(slide, sections, slide_width, slide_height)
        elif bullet_points:
            add_content_to_slide(slide, bullet_points, slide_width, slide_height)
    
    # 添加 Q&A 页
    create_qa_slide(prs)
    
    # 保存文件
    prs.save(output_file)
    print(f'成功生成 PPT 文件: {output_file}')
    print(f'总幻灯片数: {len(prs.slides)} 张')
    print()


def main():
    """主函数：处理命令行参数"""
    
    if len(sys.argv) >= 2:
        md_file = sys.argv[1]
        output_file = sys.argv[2] if len(sys.argv) >= 3 else None
    else:
        # 交互式输入
        md_file = input('请输入 Markdown 文件路径: ').strip()
        output_file = input('请输入输出 PPT 文件名 (可选，按 Enter 使用默认名): ').strip() or None
    
    try:
        create_presentation_from_markdown(md_file, output_file)
        print('=' * 60)
        print('PPT 生成成功!')
        print('=' * 60)
    except Exception as e:
        print(f'错误: {e}')
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == '__main__':
    main()
