"""
================================================================================
脚本名称: 福州中心城区公园绿地服务空间公平性评价系统 - PPT 汇报自动生成工具
语言: 中文

核心功能:
    - 内容覆盖：9张完整幻灯片，从研究背景、技术架构、指标体系、数据处理、
      到结果展示 (空间盲区、公平性分析、弱势群体) 及规划建议的完整逻辑链条。
    - 样式封装：`create_styled_slide` 函数自动生成带有专业 Header Bar 和内容卡片的幻灯片。
    - 视觉优化：设置 16:9 宽屏画布，采用微软雅黑字体，支持多种内容排版。

使用说明:
    1. 确保已安装依赖库：pip install python-pptx
    2. 直接运行脚本，将在当前目录下生成 "Fuzhou_Park_Equity_Report.pptx"。
    3. 建议配合 ArcMap、QGIS 生成的空间分析地图截图补充到幻灯片中。

作者: Fan Zhen
最后修改日期: 2026-04-14
================================================================================
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- 定义专业配色方案 ---
COLOR_TEAL = RGBColor(15, 118, 110)    # 顶部栏主色 (Teal 700)
COLOR_WHITE = RGBColor(255, 255, 255)  # 顶部栏文字色
COLOR_TEXT = RGBColor(30, 41, 59)      # 正文文字色 (Slate 800)
COLOR_ACCENT = RGBColor(59, 130, 246)  # 强调色 (Blue 500)

def create_styled_slide(prs, layout_idx, header_text):
    """创建带有专业深色 Header Bar 的幻灯片"""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # 柔和底色
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(248, 250, 252)
    bg.line.visible = False
    
    # 绘制顶部深色导航栏
    shapes = slide.shapes
    header_bar = shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.8)
    )
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = COLOR_TEAL
    header_bar.line.visible = False

    # 在 Header Bar 中添加标题
    title_box = shapes.add_textbox(Inches(0.5), Inches(0.15), prs.slide_width - Inches(1), Inches(0.5))
    tf = title_box.text_frame
    tf.text = header_text
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = COLOR_WHITE
    p.alignment = PP_ALIGN.LEFT

    return slide

def add_content_to_slide(slide, bullet_points, slide_width, slide_height):
    """为幻灯片添加项目符号列表"""
    # 内容卡片
    content_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.45),
        Inches(1.0),
        slide_width - Inches(0.9),
        slide_height - Inches(1.35),
    )
    content_card.fill.solid()
    content_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_card.line.color.rgb = RGBColor(226, 232, 240)

    body_shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.25), slide_width - Inches(1.6), slide_height - Inches(1.95))
    tf = body_shape.text_frame
    tf.word_wrap = True
    
    for i, point in enumerate(bullet_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        p.text = point
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(18)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(8)
        p.space_after = Pt(6)
        p.line_spacing = 1.3
        p.level = 0

def add_content_with_sections(slide, sections, slide_width, slide_height):
    """为幻灯片添加多个分节的内容"""
    content_card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.45),
        Inches(1.0),
        slide_width - Inches(0.9),
        slide_height - Inches(1.35),
    )
    content_card.fill.solid()
    content_card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    content_card.line.color.rgb = RGBColor(226, 232, 240)

    body_shape = slide.shapes.add_textbox(Inches(0.8), Inches(1.25), slide_width - Inches(1.6), slide_height - Inches(1.95))
    tf = body_shape.text_frame
    tf.word_wrap = True
    
    first = True
    for section in sections:
        if not first:
            tf.add_paragraph()
        first = False
        
        # 小标题
        p_title = tf.paragraphs[0] if not first else tf.add_paragraph()
        p_title.text = section.get("title", "")
        p_title.font.name = "Microsoft YaHei"
        p_title.font.size = Pt(16)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_ACCENT
        p_title.space_before = Pt(0)
        p_title.space_after = Pt(4)
        
        # 项目符号
        for point in section.get("points", []):
            p = tf.add_paragraph()
            p.text = point
            p.font.name = "Microsoft YaHei"
            p.font.size = Pt(15)
            p.font.color.rgb = COLOR_TEXT
            p.space_before = Pt(2)
            p.space_after = Pt(2)
            p.line_spacing = 1.2
            p.level = 1

def create_professional_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # --- Slide 1: 标题页 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        slide_width,
        slide_height,
    )
    title_bg.fill.solid()
    title_bg.fill.fore_color.rgb = RGBColor(241, 245, 249)
    title_bg.line.fill.background()

    accent_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0.0),
        slide_width,
        Inches(0.25),
    )
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = COLOR_TEAL
    accent_bar.line.fill.background()

    # 主标题
    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.8), Inches(1.5))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    title_p = title_tf.paragraphs[0]
    title_p.text = "福州中心城区公园绿地服务空间公平性评价系统汇报"
    title_p.font.name = "Microsoft YaHei"
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = COLOR_TEAL

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(0.9), Inches(3.7), Inches(11.8), Inches(1.2))
    subtitle_tf = subtitle_box.text_frame
    subtitle_p = subtitle_tf.paragraphs[0]
    subtitle_p.text = "基于 PostGIS 与 WebGIS 的多维评价模型研究"
    subtitle_p.font.name = "Microsoft YaHei"
    subtitle_p.font.size = Pt(24)
    subtitle_p.font.color.rgb = COLOR_ACCENT

    # 底部信息
    info_box = slide.shapes.add_textbox(Inches(0.9), Inches(5.2), Inches(11.8), Inches(1.8))
    info_tf = info_box.text_frame
    info_tf.word_wrap = True
    
    info_lines = [
        "汇报人：[您的姓名]",
        "技术栈：PostGIS + FastAPI + Leaflet + ECharts",
        "日期：2026年4月",
    ]
    for idx, line in enumerate(info_lines):
        p = info_tf.paragraphs[0] if idx == 0 else info_tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 65, 85)
        p.space_before = Pt(4)

    # --- Slide 2: 研究背景与意义 ---
    slide = create_styled_slide(prs, 1, "Slide 2: 研究背景与意义")
    add_content_to_slide(slide, [
        "政策响应：贯彻'15分钟社区生活圈'规划，探索城市公共资源分配的公平性。",
        "研究痛点：传统静态统计忽略了真实步行路网成本，难以识别微观尺度的'空间错位'。",
        "研究目标：利用 GIS 技术量化评价公园绿地的可达性，为'精准规划'提供数据支撑。",
    ], slide_width, slide_height)

    # --- Slide 3: 系统整体架构 ---
    slide = create_styled_slide(prs, 1, "Slide 3: 系统整体架构")
    
    sections_3 = [
        {
            "title": "技术路线",
            "points": [
                "数据层 (PostGIS)：存储街道空间边界、人口属性及公园数据，实现 ST_Area/ST_Intersection 空间运算",
                "逻辑层 (FastAPI)：封装公平性算法逻辑，支持多维度 API 响应",
                "表现层 (Leaflet & ECharts)：Leaflet 展示 15min 等时圈；ECharts 绘制洛伦兹曲线",
            ]
        }
    ]
    add_content_with_sections(slide, sections_3, slide_width, slide_height)

    # --- Slide 4: 核心评价指标体系 ---
    slide = create_styled_slide(prs, 1, "Slide 4: 核心评价指标体系")
    add_content_to_slide(slide, [
        "15分钟服务覆盖率：基于路网拓扑的 15min 步行范围内的可达性",
        "基尼系数 (Gini)：量化资源分配的不均衡程度 (0.0 - 1.0)",
        "洛伦兹曲线 (Lorenz)：展示人口比例与资源占有比例的偏离度",
        "老龄化倾向性：结合 elderly_rate 评估社会弱势群体的公平感知",
    ], slide_width, slide_height)

    # --- Slide 5: 数据预处理与拓扑构建 ---
    slide = create_styled_slide(prs, 1, "Slide 5: 数据预处理与拓扑构建")
    add_content_to_slide(slide, [
        "坐标统一：执行 EPSG:3857 投影坐标系转换，确保面积计算准确",
        "数据清洗：剔除人口数为 0 的无效街道，归一化 service_rate 参数",
        "动态邻域统计：支持'全区'与'街道邻域'两种动态统计模式",
    ], slide_width, slide_height)

    # --- Slide 6: 结果展示 - 空间分布与盲区 ---
    slide = create_styled_slide(prs, 1, "Slide 6: 空间分布与盲区识别")
    add_content_to_slide(slide, [
        "总体水平：全区平均服务覆盖率为 35.84%",
        "关键盲区：仓山街道、荫营镇街道为'特别加强盲区' (覆盖率 10%)",
        "空间特征：中心街道服务饱和，边缘大体量街道资源匮乏",
    ], slide_width, slide_height)

    # --- Slide 7: 结果展示 - 公平性分析 ---
    slide = create_styled_slide(prs, 1, "Slide 7: 公平性分析")
    add_content_to_slide(slide, [
        "基尼系数：实时计算结果为 0.5582，处于'高度不均衡'区间",
        "资源集中度：前 20% 的优势人口占有了 67.61% 的公园资源",
        "不均衡特征：资源严重向优势地段集中，验证了干预的紧迫性",
    ], slide_width, slide_height)

    # --- Slide 8: 结果展示 - 弱势群体可达性 ---
    slide = create_styled_slide(prs, 1, "Slide 8: 弱势群体可达性分析")
    add_content_to_slide(slide, [
        "核心矛盾：荫营镇街道老龄化比例 (32.83%) 全区最高，但覆盖率 (10%) 处于洼地",
        "群体差异：老龄化重点街道的资源配置滞后于其社会需求",
        "可达性赤字：老龄人口的服务获取力严重不足",
    ], slide_width, slide_height)

    # --- Slide 9: 结论与规划建议 ---
    slide = create_styled_slide(prs, 1, "Slide 9: 结论与规划建议")
    
    sections_9 = [
        {
            "title": "核心结论",
            "points": [
                "系统成功识别了福州中心城区的资源配置死角",
                "Gini 系数验证了干预的紧迫性",
            ]
        },
        {
            "title": "精准建议",
            "points": [
                "精准补短板：在盲区优先布局'口袋公园'",
                "适老化改造：打通慢行系统，消除空间屏障",
                "动态考核：将公平性指标纳入年度规划考核体系",
            ]
        }
    ]
    add_content_with_sections(slide, sections_9, slide_width, slide_height)

    # --- Slide 10: Q&A 页 ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    qa_bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        slide_width,
        slide_height,
    )
    qa_bg.fill.solid()
    qa_bg.fill.fore_color.rgb = RGBColor(248, 250, 252)
    qa_bg.line.fill.background()

    qa_accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0),
        Inches(0),
        slide_width,
        Inches(0.25),
    )
    qa_accent.fill.solid()
    qa_accent.fill.fore_color.rgb = COLOR_TEAL
    qa_accent.line.fill.background()

    # Q&A 内容
    qa_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(3))
    qa_tf = qa_box.text_frame
    qa_tf.word_wrap = True
    
    qa_lines = ["Q & A", "感谢聆听，欢迎提问！"]
    for idx, line in enumerate(qa_lines):
        p = qa_tf.paragraphs[0] if idx == 0 else qa_tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(54 if idx == 0 else 32)
        p.font.bold = True if idx == 0 else False
        p.font.color.rgb = COLOR_TEAL if idx == 0 else RGBColor(71, 85, 105)
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(0)
        p.space_after = Pt(8)

    # 保存文件
    file_name = "Fuzhou_Park_Equity_Report.pptx"
    prs.save(file_name)
    print(f"✓ 成功生成 PPT 文件: {file_name}")
    print(f"✓ 总幻灯片数: 10 张")
    print(f"✓ 包含: 标题页 + 9个内容页 + Q&A页")
    print(f"✓ 配色方案: Teal + Slate 专业版")

if __name__ == "__main__":
    create_professional_presentation()
